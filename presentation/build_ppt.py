"""Build a highly premium, clean light-themed 15-slide PowerPoint presentation with custom native flowcharts and block layouts."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(__file__), "output", "IdeaRound_Enterprise_Knowledge_Copilot.pptx")

# ── Color Palette (Clean, Professional, Light Theme) ──────────────────────────
BG_COLOR = RGBColor(0xF8, 0xFA, 0xFC)        # Light Slate/White background
TITLE_COLOR = RGBColor(0x1E, 0x2E, 0x4A)     # Deep Navy Blue
ACCENT_COLOR = RGBColor(0x2B, 0x6C, 0xE2)    # Royal Blue (Category Pill & Highlights)
TEXT_COLOR = RGBColor(0x33, 0x41, 0x55)      # Dark Charcoal
MUTED_COLOR = RGBColor(0x71, 0x80, 0x96)     # Cool Slate
HIGHLIGHT_COLOR = RGBColor(0xC5, 0x30, 0x30) # Warm Crimson (for bottlenecks/metrics)
GREEN_COLOR = RGBColor(0x2F, 0x85, 0x5A)     # Deep Green (for wins/speedups)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)         # Solid White for cards
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)     # Light Slate border
BLOCK_BORDER = CARD_BORDER                   # Alias for connector backgrounds

# ── Category Colors ──────────────────────────────────────────────────────────
PURPLE_CAT = RGBColor(0x80, 0x5A, 0xD5)
BLUE_CAT = RGBColor(0x2B, 0x6C, 0xE2)
CYAN_CAT = RGBColor(0x00, 0xA3, 0xC4)
GREEN_CAT = RGBColor(0x2F, 0x85, 0x5A)
ORANGE_CAT = RGBColor(0xDD, 0x6B, 0x20)
RED_CAT = RGBColor(0xC5, 0x30, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_category_pill(slide, text):
    width = 0.55 + len(text) * 0.085
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.35), Inches(width), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background() # borderless
    
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.72), Inches(12.0), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_top = Inches(0)
    p_title = tf_title.paragraphs[0]
    p_title.text = text
    p_title.font.size = Pt(30)
    p_title.font.bold = True
    p_title.font.color.rgb = TITLE_COLOR
    p_title.font.name = "Arial"

def draw_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    return card

def draw_flow_step(slide, left, top, width, height, text, bg_color, symbol=""):
    card = draw_card(slide, left, top, width, height, bg_color=bg_color, border_color=None)
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(width), Inches(height - 0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    
    p = tf.paragraphs[0]
    p.text = symbol + " " + text if symbol else text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"

def draw_circle_icon(slide, left, top, size, bg_color, symbol):
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    oval.fill.solid()
    oval.fill.fore_color.rgb = bg_color
    oval.line.fill.background()
    
    tf = oval.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = symbol
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"

def add_card_text(slide, left, top, width, height, title, body_lines, title_color=TITLE_COLOR, title_size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Arial"
    
    for line in body_lines:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_COLOR
        p2.font.name = "Calibri"
        p2.space_before = Pt(4)
    return tf

def draw_connector_line(slide, x1, y1, x2, y2, color=ACCENT_COLOR, width=1.5):
    shape = slide.shapes.add_shape(1, Inches(x1), Inches(y1), Inches(x2 - x1), Inches(y2 - y1) if (y2 - y1) > 0.01 else Pt(width))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_bottom_bar(slide, text):
    draw_card(slide, 0.6, 6.2, 12.133, 0.5, bg_color=RGBColor(0xEB, 0xF8, 0xFF), border_color=None)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12.133), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = "Arial"

# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)

# Center Card
draw_card(s, 1.5, 1.5, 10.333, 4.5)

# Category pill inside title card
width = 0.55 + len("PROJECT OVERVIEW") * 0.085
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.666 - width/2), Inches(2.0), Inches(width), Inches(0.32))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_COLOR
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "PROJECT OVERVIEW"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(8.5)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = "Arial"

# Main Project Title
title_box = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🛡️ Enterprise Knowledge Copilot"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.font.name = "Arial"

p2 = tf.add_paragraph()
p2.text = "A Privacy-First, Hybrid Retrieval RAG Assistant for Secure IT Support"
p2.alignment = PP_ALIGN.CENTER
p2.font.size = Pt(16)
p2.font.color.rgb = MUTED_COLOR
p2.font.name = "Arial"
p2.space_before = Pt(8)

# Author Info Card
info_box = s.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.2))
tf_i = info_box.text_frame
tf_i.word_wrap = True
p_name = tf_i.paragraphs[0]
p_name.text = "Ichhit Karwa"
p_name.alignment = PP_ALIGN.CENTER
p_name.font.size = Pt(18)
p_name.font.bold = True
p_name.font.color.rgb = TITLE_COLOR
p_name.font.name = "Arial"

p_uni = tf_i.add_paragraph()
p_uni.text = "VIT Bhopal University  ·  B.Tech Computer Science  ·  3rd Year"
p_uni.alignment = PP_ALIGN.CENTER
p_uni.font.size = Pt(13)
p_uni.font.color.rgb = MUTED_COLOR
p_uni.font.name = "Arial"
p_uni.space_before = Pt(4)

p_event = tf_i.add_paragraph()
p_event.text = "IdeaRound Submission  ·  Technical Presentation"
p_event.alignment = PP_ALIGN.CENTER
p_event.font.size = Pt(11)
p_event.font.color.rgb = MUTED_COLOR
p_event.font.name = "Arial"
p_event.space_before = Pt(16)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: THE CHALLENGE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "THE CHALLENGE")
add_title(s, "Enterprise Knowledge Fragmentation")

# Left Card: Gaps
draw_card(s, 0.6, 1.65, 5.8, 4.9)
add_card_text(s, 0.9, 1.95, 5.2, 4.3, "❌ Gaps in Legacy Support Systems", [
    "• Scattered Document Corpora:",
    "  Support staff spend hours searching across static wikis, PDF runbooks, and historical system handbooks.",
    "",
    "• Ingestion Scaling Pain Points:",
    "  Traditional databases scale poorly, making queries across hundreds of thousands of historical tickets slow.",
    "",
    "• Ticket Duplication Swarms:",
    "  Lack of immediate historical search leads to users submitting duplicate tickets for known resolved issues."
], title_size=16)

# Right Card: Privacy Constraints
draw_card(s, 6.8, 1.65, 5.9, 4.9)
add_card_text(s, 7.1, 1.95, 5.3, 4.3, "🔒 Privacy & PII Constraints", [
    "• Compliance Restrictions:",
    "  Enterprise guidelines forbid sharing names, phone numbers, email addresses, and internal system IPs.",
    "",
    "• Public SaaS Vulnerability:",
    "  Sending internal server logs to public APIs (OpenAI, Anthropic) risks data-leakage and compliance failures.",
    "",
    "• Ingestion-Level Mitigation:",
    "  Local ingestion scripts redact PII (using SpaCy NER + regex scrubs) before storing documents or vectors locally."
], title_size=16)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: WHY EXISTING SEARCH FAILS (NEW!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "PROBLEM VS SOLUTION")
add_title(s, "Why Existing Search Fails")

# Table
table_left = 0.6
table_top = 1.8
table_width = 12.133
table_height = 4.0

rows = 5
cols = 3
table_shape = s.shapes.add_table(rows, cols, Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_height))
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(4.316)
table.columns[2].width = Inches(4.316)

headers = ["User Query", "Traditional Search", "Hybrid Copilot"]
data = [
    ["TKT-10084", "Exact match only", "Instant ticket lookup"],
    ["VPN not working after update", "Poor recall", "Semantic retrieval"],
    ["Explain VPN issue in Hindi", "Impossible", "LLM grounded response"],
    ["Duplicate ticket creation", "No detection", "Similarity check"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = TITLE_COLOR
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(13)

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT_COLOR
            if col_idx == 2:
                paragraph.font.bold = True
                paragraph.font.color.rgb = GREEN_COLOR

add_bottom_bar(s, "Traditional queries fail under semantic variation and cross-lingual needs. Hybrid Copilot handles them natively.")


# ═══════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE OVERVIEW (NEW - HERO FLOW!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "ARCHITECTURE")
add_title(s, "System Architecture Flow")

# Left Column (Hero Description)
draw_card(s, 0.6, 1.65, 5.0, 4.9)
add_card_text(s, 0.9, 2.0, 4.4, 4.3, "⚙️ Dynamic Request Execution Flow", [
    "• Single Point of Entry:",
    "  The user query initiates a state context within FastAPI or Streamlit UI.",
    "",
    "• LangGraph State Machine Routing:",
    "  Intently parses parameters and calls specific tools based on dynamic decision nodes.",
    "",
    "• Database Offload:",
    "  Vector search (ChromaDB) and token keyword matching (BM25) run in parallel, resolving queries against local memory structures.",
    "",
    "• Safe Generative Grounding:",
    "  The local LLaMA model receives redacted context to build the final answer securely."
], title_size=16)

# Right Column: Big Vertical Flowchart
# Draw vertical connectors behind blocks
draw_connector_line(s, 9.45, 1.6, 9.45, 6.3, color=BLOCK_BORDER, width=2.5)

flow_steps = [
    ("User", RED_CAT, 1.65),
    ("FastAPI / Streamlit", ORANGE_CAT, 2.3),
    ("LangGraph Agent", PURPLE_CAT, 2.95),
    ("Tools", BLUE_CAT, 3.6),
    ("Hybrid Retrieval", CYAN_CAT, 4.25),
    ("ChromaDB / BM25", GREEN_CAT, 4.9),
    ("LLM", ORANGE_CAT, 5.55),
    ("Answer", GREEN_CAT, 6.2)
]


for text, color, top in flow_steps:
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(top), Inches(5.3), Inches(0.45))
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()
    p = block.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"


# ═══════════════════════════════════════════════════════════
# SLIDE 5: 6-LAYER SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "ARCHITECTURE")
add_title(s, "6-Layer System Architecture")

# Draw vertical connecting line behind squares
draw_connector_line(s, 0.875, 1.8, 0.875, 5.8, color=ACCENT_COLOR, width=2)

layers = [
    (6, "Presentation Layer", "Streamlit Chat UI · Health Dashboard · Feedback", PURPLE_CAT),
    (5, "Orchestration Layer", "LangGraph State Machine · Router / Planner · 6 Tools", BLUE_CAT),
    (4, "Cognitive Layer", "LLaMA 3.2 (3B) via Ollama · 100% Local Inference", CYAN_CAT),
    (3, "Retrieval Layer", "Hybrid Search: ChromaDB + BM25 + NetworkX → RRF → Cross-Encoder", GREEN_CAT),
    (2, "Data & Security Layer", "spaCy PII Redaction · Embedding · Elasticsearch", ORANGE_CAT),
    (1, "Infrastructure Layer", "Ollama Server · FastAPI · Streamlit · Local Storage", RED_CAT)
]

for idx, (num, title, details, color) in enumerate(layers):
    top = 1.65 + idx * 0.8
    # Draw number square
    draw_card(s, 0.6, top, 0.55, 0.55, bg_color=color, border_color=None)
    # Put number text
    tx = s.shapes.add_textbox(Inches(0.6), Inches(top + 0.05), Inches(0.55), Inches(0.45))
    p = tx.text_frame.paragraphs[0]
    p.text = str(num)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"
    
    # Draw white card on the right
    draw_card(s, 1.3, top, 11.4, 0.55)
    # Add text inside the card
    tx_card = s.shapes.add_textbox(Inches(1.5), Inches(top + 0.05), Inches(11.0), Inches(0.45))
    tf = tx_card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(12.5)
    p_title.font.bold = True
    p_title.font.color.rgb = TITLE_COLOR
    p_title.font.name = "Arial"
    
    p_det = tf.add_paragraph()
    p_det.text = details
    p_det.font.size = Pt(9.5)
    p_det.font.color.rgb = MUTED_COLOR
    p_det.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════
# SLIDE 6: WHY HYBRID? (NEW!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "RETRIEVAL PIPELINE")
add_title(s, "Why Hybrid Search?")

# Top query block
draw_card(s, 4.666, 1.65, 4.0, 0.6, bg_color=TITLE_COLOR, border_color=None)
tx = s.shapes.add_textbox(Inches(4.666), Inches(1.7), Inches(4.0), Inches(0.5))
p = tx.text_frame.paragraphs[0]
p.text = "User Query"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(255, 255, 255)

# Connectors
draw_connector_line(s, 6.666, 2.25, 6.666, 2.5, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 3.666, 2.5, 9.666, 2.5, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 3.666, 2.5, 3.666, 2.95, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 9.666, 2.5, 9.666, 2.95, color=ACCENT_COLOR, width=1.5)

# Two retrieval paths
draw_card(s, 1.666, 2.95, 4.0, 0.7, bg_color=PURPLE_CAT, border_color=None)
tx1 = s.shapes.add_textbox(Inches(1.666), Inches(3.05), Inches(4.0), Inches(0.5))
p = tx1.text_frame.paragraphs[0]
p.text = "BM25 Keyword Search\n(Exact tokens / SKU codes)"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(255, 255, 255)

draw_card(s, 7.666, 2.95, 4.0, 0.7, bg_color=BLUE_CAT, border_color=None)
tx2 = s.shapes.add_textbox(Inches(7.666), Inches(3.05), Inches(4.0), Inches(0.5))
p = tx2.text_frame.paragraphs[0]
p.text = "Vector Search\n(Semantic concepts)"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(255, 255, 255)

# Connectors back to fusion
draw_connector_line(s, 3.666, 3.65, 3.666, 4.0, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 9.666, 3.65, 9.666, 4.0, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 3.666, 4.0, 9.666, 4.0, color=ACCENT_COLOR, width=1.5)
draw_connector_line(s, 6.666, 4.0, 6.666, 4.25, color=ACCENT_COLOR, width=1.5)

# Fusion Block
draw_card(s, 4.666, 4.25, 4.0, 0.55, bg_color=CYAN_CAT, border_color=None)
tx3 = s.shapes.add_textbox(Inches(4.666), Inches(4.3), Inches(4.0), Inches(0.5))
p = tx3.text_frame.paragraphs[0]
p.text = "RRF Fusion"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(255, 255, 255)

draw_connector_line(s, 6.666, 4.8, 6.666, 5.05, color=ACCENT_COLOR, width=1.5)

# Cross Encoder
draw_card(s, 4.666, 5.05, 4.0, 0.55, bg_color=GREEN_CAT, border_color=None)
tx4 = s.shapes.add_textbox(Inches(4.666), Inches(5.1), Inches(4.0), Inches(0.5))
p = tx4.text_frame.paragraphs[0]
p.text = "Cross Encoder Reranker"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(255, 255, 255)

draw_connector_line(s, 6.666, 5.6, 6.666, 5.85, color=ACCENT_COLOR, width=1.5)

# Top 5 Chunks / LLM
draw_card(s, 4.666, 5.85, 4.0, 0.55, bg_color=ORANGE_CAT, border_color=None)
tx5 = s.shapes.add_textbox(Inches(4.666), Inches(5.9), Inches(4.0), Inches(0.5))
p = tx5.text_frame.paragraphs[0]
p.text = "Top 5 Chunks passed to LLM"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(11.5)
p.font.color.rgb = RGBColor(255, 255, 255)

add_bottom_bar(s, "BM25 handles exact terms while vector search captures semantic meaning.")


# ═══════════════════════════════════════════════════════════
# SLIDE 7: WHY ELASTICSEARCH? (NEW!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "SCALING")
add_title(s, "Why Elasticsearch?")

# Table
table_left = 0.6
table_top = 1.8
table_width = 12.133
table_height = 4.0

rows = 6
cols = 2
table_shape = s.shapes.add_table(rows, cols, Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_height))
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(6.066)
table.columns[1].width = Inches(6.066)

headers = ["Local Stack", "Elasticsearch"]
data = [
    ["BM25 Python (rank-bm25 library)", "Lucene Index (native index structures)"],
    ["Chroma (in-process database)", "kNN Search (distributed vector database)"],
    ["Single Process (Python GIL bounds)", "Concurrent Search (distributed cluster worker pools)"],
    ["Higher contention under concurrent read/write", "Stable throughput via concurrent segment search"],
    ["Separate systems (Chroma + rank-bm25 python)", "Unified engine (combined text + vector search)"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = TITLE_COLOR
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(13)

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT_COLOR
            if col_idx == 1:
                paragraph.font.bold = True
                paragraph.font.color.rgb = ACCENT_COLOR

add_bottom_bar(s, "Elasticsearch consolidates sparse keyword matching and dense vector search into a single, scalable database.")


# ═══════════════════════════════════════════════════════════
# SLIDE 8: RAG ARCHITECTURE SHIFT (NEW!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "PERFORMANCE")
add_title(s, "Local Python RAG vs. Dedicated Search Engine")

# Track 1: In-Process RAG
draw_flow_step(s, 0.6, 2.2, 2.0, 1.0, "In-Process RAG\n(Chroma + BM25)", TITLE_COLOR)

# Connectors for Track 1
draw_connector_line(s, 2.6, 2.7, 2.8, 2.7, color=MUTED_COLOR, width=2)
draw_connector_line(s, 4.5, 2.7, 4.8, 2.7, color=MUTED_COLOR, width=2)
draw_connector_line(s, 6.5, 2.7, 6.8, 2.7, color=MUTED_COLOR, width=2)
draw_connector_line(s, 8.5, 2.7, 8.8, 2.7, color=MUTED_COLOR, width=2)
draw_connector_line(s, 10.5, 2.7, 10.8, 2.7, color=MUTED_COLOR, width=2)

draw_flow_step(s, 2.8, 2.2, 1.7, 1.0, "50 Users\n(Concurrent)", MUTED_COLOR, "👥")
draw_flow_step(s, 4.8, 2.2, 1.7, 1.0, "Single Python\nThread (GIL)", ORANGE_CAT, "🧵")
draw_flow_step(s, 6.8, 2.2, 1.7, 1.0, "On-the-fly Scan\n(200k docs)", RED_CAT, "🔄")
draw_flow_step(s, 8.8, 2.2, 1.7, 1.0, "Memory Bloat\n(1.5 GB RAM)", HIGHLIGHT_COLOR, "💾")
draw_flow_step(s, 10.8, 2.2, 1.7, 1.0, "23.0s Latency\n(Queue Build-up)", HIGHLIGHT_COLOR, "💀")

# Track 2: Elasticsearch RAG
draw_flow_step(s, 0.6, 4.2, 2.0, 1.0, "Engine RAG\n(Elasticsearch)", ACCENT_COLOR)

# Connectors for Track 2
draw_connector_line(s, 2.6, 4.7, 2.8, 4.7, color=ACCENT_COLOR, width=2)
draw_connector_line(s, 4.5, 4.7, 4.8, 4.7, color=ACCENT_COLOR, width=2)
draw_connector_line(s, 6.5, 4.7, 6.8, 4.7, color=ACCENT_COLOR, width=2)
draw_connector_line(s, 8.5, 4.7, 8.8, 4.7, color=ACCENT_COLOR, width=2)
draw_connector_line(s, 10.5, 4.7, 10.8, 4.7, color=ACCENT_COLOR, width=2)

draw_flow_step(s, 2.8, 4.2, 1.7, 1.0, "50 Users\n(Concurrent)", MUTED_COLOR, "👥")
draw_flow_step(s, 4.8, 4.2, 1.7, 1.0, "Elasticsearch\nThread Pool", BLUE_CAT, "⚙️")
draw_flow_step(s, 6.8, 4.2, 1.7, 1.0, "Pre-built Lucene\nIndex Lookup", GREEN_CAT, "📖")
draw_flow_step(s, 8.8, 4.2, 1.7, 1.0, "Compressed RAM\n+ OS Disk Cache", CYAN_CAT, "⚡")
draw_flow_step(s, 10.8, 4.2, 1.7, 1.0, "1.7s Latency\n(Concurrent)", GREEN_COLOR, "🚀")

# Bottom explanations (subtitles under the tracks)
txBox_left = s.shapes.add_textbox(Inches(2.8), Inches(3.35), Inches(9.7), Inches(0.6))
tf_l = txBox_left.text_frame
tf_l.word_wrap = True
p_l = tf_l.paragraphs[0]
p_l.text = "⚠️ Bottleneck: Python performs search calculations during query time, causing massive CPU & memory overhead."
p_l.font.size = Pt(11)
p_l.font.bold = True
p_l.font.color.rgb = HIGHLIGHT_COLOR
p_l.font.name = "Calibri"

txBox_right = s.shapes.add_textbox(Inches(2.8), Inches(5.35), Inches(9.7), Inches(0.6))
tf_r = txBox_right.text_frame
tf_r.word_wrap = True
p_r = tf_r.paragraphs[0]
p_r.text = "✅ Advantage: Elasticsearch does heavy indexing work beforehand. Query time search is just a fast precomputed lookup."
p_r.font.size = Pt(11)
p_r.font.bold = True
p_r.font.color.rgb = GREEN_COLOR
p_r.font.name = "Calibri"

add_bottom_bar(s, "One Sentence for Jury: Elasticsearch computes much less during query time because the work is done during indexing.")


# ═══════════════════════════════════════════════════════════
# SLIDE 9: LANGGRAPH AGENT & TOOLS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "AUTONOMOUS AGENT")
add_title(s, "LangGraph State Machine — 6 Specialised Tools")

# Center Top: Router / Planner Box
draw_card(s, 4.166, 1.65, 5.0, 0.9, bg_color=TITLE_COLOR, border_color=None)
tx = s.shapes.add_textbox(Inches(4.166), Inches(1.75), Inches(5.0), Inches(0.7))
p = tx.text_frame.paragraphs[0]
p.text = "Router / Planner\nDecides best action based on query intent"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = "Arial"

# Draw flow lines from Planner box to Columns
draw_connector_line(s, 6.666, 2.55, 6.666, 2.9, color=ACCENT_COLOR, width=1.5) # center line
draw_connector_line(s, 2.5, 2.75, 10.833, 2.75, color=ACCENT_COLOR, width=1.5) # horizontal bar
draw_connector_line(s, 2.5, 2.75, 2.5, 2.9, color=ACCENT_COLOR, width=1.5) # left vertical drop
draw_connector_line(s, 10.833, 2.75, 10.833, 2.9, color=ACCENT_COLOR, width=1.5) # right vertical drop

tools_grid = [
    # (col, row, title, detail, color, symbol)
    (0, 0, "tool_search_docs", ["Policies, runbooks, handbook"], PURPLE_CAT, "🔍"),
    (1, 0, "tool_search_tickets", ["Historical incident reports"], BLUE_CAT, "🗄️"),
    (2, 0, "tool_filtered_tickets", ["Metadata filter: P1, VPN, Auth"], CYAN_CAT, "🥞"),
    (0, 1, "tool_summarize", ["High-level overview across all data"], GREEN_CAT, "📊"),
    (1, 1, "tool_multihop", ["Cross-source reasoning: docs + graph"], ORANGE_CAT, "🔗"),
    (2, 1, "tool_create_ticket", ["Creates ticket + dedup check (>0.75)"], RED_CAT, "✔")
]

for col, row, title, detail, color, symbol in tools_grid:
    left = 0.6 + col * 4.166
    top = 2.9 + row * 1.65
    draw_card(s, left, top, 3.8, 1.4)
    draw_circle_icon(s, left + 0.25, top + 0.35, 0.7, color, symbol)
    add_card_text(s, left + 1.1, top + 0.3, 2.6, 1.0, title, detail, title_color=TITLE_COLOR, title_size=13.5)

# Bottom Bar
add_bottom_bar(s, "✓ Memory (Conversation + Entity)        ✓ PII Guardrails + Safety        ✓ Retry & Fallbacks        ✓ Audit Logging → audit.jsonl")


# ═══════════════════════════════════════════════════════════
# SLIDE 10: ENGINEERING CHALLENGES & DEBUGGING STORIES
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "ENGINEERING")
add_title(s, "Engineering Challenges & Resolutions")

# Left Card: Memory Explosion
draw_card(s, 0.6, 1.65, 5.8, 4.9)
add_card_text(s, 0.9, 2.0, 5.2, 4.3, "💾 1. Memory Explosion", [
    "• Initial State:",
    "  Spaced out ingestion models holding massive structures spiked local memory usage up to 7GB RAM.",
    "",
    "• Debugging & Resolution Flow:",
    "  Memory Spike (7GB RAM)",
    "  ↓",
    "  Cache optimization & lazy evaluation",
    "  ↓",
    "  Stable ingestion memory (<5GB RAM)"
], title_size=16)

# Right Card: JVM Heap
draw_card(s, 6.8, 1.65, 5.9, 4.9)
add_card_text(s, 7.1, 2.0, 5.3, 4.3, "⚡ 2. JVM Heap Constraints", [
    "• Initial State:",
    "  Elasticsearch bulk indexing of 200,000 documents crashed due to container limit defaults.",
    "",
    "• Debugging & Resolution Flow:",
    "  JVM Heap OOM (512MB RAM)",
    "  ↓",
    "  Raised settings to 2GB",
    "  ↓",
    "  Stable continuous indexing"
], title_size=16)


# ═══════════════════════════════════════════════════════════
# SLIDE 11: BENCHMARKS (SAFE WORKLOADS)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "BENCHMARKS")
add_title(s, "Concurrent Semantic Search (50 Users)")

# Create Table
table_left = 0.6
table_top = 1.8
table_width = 6.2
table_height = 3.5

rows = 4
cols = 3
table_shape = s.shapes.add_table(rows, cols, Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_height))
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)

headers = ["Metric", "Local", "Elasticsearch"]
data = [
    ["P50 Latency", "23s", "1.7s"],
    ["P95 Latency", "41s", "1.8s"],
    ["Throughput", "1.27 RPS", "28.9 RPS"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = TITLE_COLOR
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(12)

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11.5)
            paragraph.font.color.rgb = TEXT_COLOR
            if col_idx == 2:
                paragraph.font.bold = True
                paragraph.font.color.rgb = GREEN_COLOR

# Left explanation below table
txBox = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(6.2), Inches(1.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "▸ Analysis:\nLocal execution throttles under multi-threaded queue build-up. Elasticsearch utilizes multi-threaded search across segment blocks to maintain low latency."
p.font.size = Pt(11)
p.font.color.rgb = TEXT_COLOR
p.font.name = "Calibri"

# Right Card
draw_card(s, 7.2, 1.8, 5.5, 4.2)
add_card_text(s, 7.5, 2.1, 4.9, 3.6, "📈 Performance Conclusion", [
    "• Stable Latency under Load:",
    "  Under concurrent workloads, Elasticsearch maintained much more stable latency.",
    "",
    "• Native Threading Advantage:",
    "  Lucene segment readers process vector arithmetic and token matches in parallel, preventing single-threaded locks.",
    "",
    "• Latency Reductions:",
    "  P95 tails remain bound close to average lookup speeds, preventing worst-case query freezes."
], title_size=15)


# ═══════════════════════════════════════════════════════════
# SLIDE 12: KEY LEARNINGS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "KEY LEARNINGS")
add_title(s, "Key Engineering Learnings")

learnings = [
    (0, 0, "1. Profiling > Assumptions", ["Assumptions about code bottlenecks are usually wrong. Run profiling to find the exact CPU/memory hogs."], PURPLE_CAT, "📊"),
    (1, 0, "2. Hybrid Improves Quality", ["Combining token BM25 with vector semantics results in significantly better search recall."], BLUE_CAT, "🧬"),
    (2, 0, "3. Bypass Vector Search", ["Exact ticket ID queries should bypass vector/LLM cycles entirely for instant <1ms lookup."], CYAN_CAT, "⚡"),
    (0, 1, "4. Inference Dominates", ["Local LLM model inference takes up 80%+ of the response latency. Retrieval is only a fraction."], GREEN_CAT, "🧠"),
    (1, 1, "5. Caching is Critical", ["In-memory caching of org charts and frequent queries eliminates redundant DB operations."], ORANGE_CAT, "💾"),
    (2, 1, "6. Concurrency Bottlenecks", ["What runs fast in a single-user shell will fail under simultaneous concurrent requests."], RED_CAT, "🔀")
]

for col, row, title, detail, color, symbol in learnings:
    left = 0.6 + col * 4.166
    top = 1.9 + row * 2.1
    draw_card(s, left, top, 3.8, 1.8)
    draw_circle_icon(s, left + 0.25, top + 0.25, 0.55, color, symbol)
    add_card_text(s, left + 0.95, top + 0.2, 2.7, 1.4, title, detail, title_color=TITLE_COLOR, title_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 13: LIMITATIONS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "LIMITATIONS")
add_title(s, "Current Limitations")

# Left Card: Infrastructure constraints
draw_card(s, 0.6, 1.65, 5.8, 4.9)
add_card_text(s, 0.9, 2.0, 5.2, 4.3, "⚠️ Infrastructure constraints", [
    "• Single-Node Deployment:",
    "  The current prototype runs Elasticsearch and the FastAPI server on a single machine, lacking high-availability redundancy.",
    "",
    "• Local Ollama Inference Latency:",
    "  Generating responses using local models is bounded by host GPU memory speeds, resulting in 2-3s token completion latency."
], title_size=16)

# Right Card: Scale & Cache limits
draw_card(s, 6.8, 1.65, 5.9, 4.9)
add_card_text(s, 7.1, 2.0, 5.3, 4.3, "🔒 Prototype & Cache Limits", [
    "• Retrieval Quality ceiling:",
    "  Needs custom domain embedding weights fine-tuning to better understand specific corporate IT terminology.",
    "",
    "• No Redis Caching Layer:",
    "  Duplicate query lookups hit the database repeatedly, wasting CPU cycles on identical retrievals.",
    "",
    "• Prototype Scale limits:",
    "  Currently validated up to 200,000 ticket documents. Production demands scaling to millions."
], title_size=16)


# ═══════════════════════════════════════════════════════════
# SLIDE 14: ROADMAP
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "ROADMAP")
add_title(s, "Future Production Scaling Roadmap")

# Center Top: Roadmap Header Box
draw_card(s, 4.166, 1.65, 5.0, 0.9, bg_color=TITLE_COLOR, border_color=None)
tx = s.shapes.add_textbox(Inches(4.166), Inches(1.75), Inches(5.0), Inches(0.7))
p = tx.text_frame.paragraphs[0]
p.text = "Architecture supports future horizontal scaling"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(12.5)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = "Arial"

# Horizontal Flowchart
# Draw connecting line behind blocks
draw_connector_line(s, 1.4, 3.8, 11.9, 3.8, color=BLOCK_BORDER, width=2)

roadmap_steps = [
    ("React", PURPLE_CAT, 0.6),
    ("FastAPI", BLUE_CAT, 2.5),
    ("Redis", CYAN_CAT, 4.4),
    ("Elasticsearch \n Cluster", GREEN_CAT, 6.3),
    ("vLLM", ORANGE_CAT, 8.2),
    ("AWS EKS", RED_CAT, 10.1)
]

for text, color, left in roadmap_steps:
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.2), Inches(1.8), Inches(1.2))
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()
    
    tf = block.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Arial"

add_bottom_bar(s, "✓ UI Integration (React)        ✓ Redis Caching        ✓ Distributed ES shards        ✓ vLLM GPU Server (AWS EKS)")


# ═══════════════════════════════════════════════════════════
# SLIDE 15: LIVE DEMO SCENARIOS
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s)
add_category_pill(s, "LIVE DEMO")
add_title(s, "Interactive System Scenarios")

scenarios = [
    ("🎫 Exact Lookup", "User Query:", "TKT-10084", "System Behavior:", "• Direct query routing bypasses the vector search and LLM entirely.\n• Pulls matching status and logs in <1ms.", PURPLE_CAT),
    ("🔎 Semantic Search", "User Query:", "Why VPN failing?", "System Behavior:", "• Hybrid retrieval fetches tickets with similar network issue tokens.\n• Summarizes logs into a concise response.", BLUE_CAT),
    ("🌐 Hindi Support", "User Query:", "Explain VPN in Hindi", "System Behavior:", "• Resolves retrieval in English and leverages LLM capabilities to formulate response in Hindi.", CYAN_CAT),
    ("☠️ Pirate Mode", "User Query:", "Explain like a pirate", "System Behavior:", "• Demonstrates instruction flexibility in local generation models to alter response tone.", ORANGE_CAT)
]

for idx, (title, q_lbl, query, s_lbl, system, color) in enumerate(scenarios):
    left = 0.6 + idx * 3.033
    draw_card(s, left, 1.8, 2.9, 4.1)
    
    # Text Box
    txBox = s.shapes.add_textbox(Inches(left + 0.15), Inches(1.95), Inches(2.6), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    
    # Header
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Arial"
    
    # Query Label
    p2 = tf.add_paragraph()
    p2.text = q_lbl
    p2.font.size = Pt(10)
    p2.font.bold = True
    p2.font.color.rgb = MUTED_COLOR
    p2.space_before = Pt(8)
    
    # Query Value
    p3 = tf.add_paragraph()
    p3.text = f'"{query}"'
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = TITLE_COLOR
    
    # System Label
    p4 = tf.add_paragraph()
    p4.text = s_lbl
    p4.font.size = Pt(10)
    p4.font.bold = True
    p4.font.color.rgb = MUTED_COLOR
    p4.space_before = Pt(8)
    
    # System behavior bullets
    for line in system.split("\n"):
        p5 = tf.add_paragraph()
        p5.text = line
        p5.font.size = Pt(10)
        p5.font.color.rgb = TEXT_COLOR
        p5.font.name = "Calibri"

add_bottom_bar(s, "Interactive UI switches adapt routing rules instantly to demo custom LLM instructions and direct bypass pipelines.")

prs.save(OUT)
print(f"DONE: {OUT}")
